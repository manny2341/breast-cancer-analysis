"""
Creates a PowerPoint presentation for the Breast Cancer Detection project
Simple English — easy to understand
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json, os

# ── COLOURS ──────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1a, 0x1a, 0x2e)
RED        = RGBColor(0xc0, 0x39, 0x2b)
GREEN      = RGBColor(0x27, 0xae, 0x60)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW     = RGBColor(0xFF, 0xCC, 0x00)
LIGHT_GREY = RGBColor(0xAA, 0xAA, 0xAA)
PINK       = RGBColor(0xFF, 0xCC, 0xCC)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank layout


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color=DARK_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, text, x, y, w, h,
        font_size=18, bold=False, color=WHITE,
        bg_color=None, align=PP_ALIGN.LEFT,
        italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf    = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic    = italic
    if bg_color:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txBox


def rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y),
                                 Inches(w), Inches(h))


# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 7.5, DARK_BG)
rect(sl, 0, 2.8, 13.33, 2.0, RED)

box(sl, "🎗️ Breast Cancer Detection System",
    0.5, 2.9, 12, 1.0, font_size=36, bold=True,
    color=WHITE, align=PP_ALIGN.CENTER)

box(sl, "Using Machine Learning to Detect Cancer from Cell Measurements",
    0.5, 3.9, 12, 0.6, font_size=20,
    color=PINK, align=PP_ALIGN.CENTER)

box(sl, "569 Patients  ·  6 Machine Learning Models  ·  97.5% Accuracy",
    0.5, 4.6, 12, 0.5, font_size=16,
    color=YELLOW, align=PP_ALIGN.CENTER)

box(sl, "COS5029-B: Data Science for AI",
    0.5, 6.5, 12, 0.5, font_size=14,
    color=LIGHT_GREY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT IS THIS PROJECT?
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "What Is This Project?", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

box(sl, "The Problem",
    0.5, 1.2, 5.5, 0.5, font_size=20, bold=True, color=YELLOW)
box(sl, "When a doctor finds a lump, they take a tiny sample of cells.\n"
        "A computer then measures the shape and size of those cells.\n\n"
        "The question is:\n"
        "Is this lump CANCEROUS or SAFE?",
    0.5, 1.8, 5.5, 2.5, font_size=16, color=WHITE)

box(sl, "Our Solution",
    7.3, 1.2, 5.5, 0.5, font_size=20, bold=True, color=YELLOW)
box(sl, "We trained a computer to look at those measurements\n"
        "and answer that question automatically.\n\n"
        "It learned from 569 real patients\n"
        "and can now detect cancer with 97.5% accuracy.",
    7.3, 1.8, 5.5, 2.5, font_size=16, color=WHITE)

rect(sl, 0.5, 4.5, 12.3, 1.8, RGBColor(0x2a, 0x2a, 0x4e))
box(sl, "Think of it like this:",
    0.8, 4.6, 11, 0.4, font_size=14, bold=True, color=YELLOW)
box(sl, "A very experienced doctor has seen thousands of cancer cases."
        " Over time they learn what cancerous cells look like."
        " Our model does the same thing — it studied 569 cases and learned the patterns.",
    0.8, 5.0, 11.5, 1.1, font_size=14, color=LIGHT_GREY, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 3 — THE DATA
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "The Data — What Did We Work With?", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

# Stats boxes
for x, val, label, col in [
    (0.5,  "569",   "Total Patients",          GREEN),
    (3.5,  "357",   "Benign\n(No Cancer)",      GREEN),
    (6.5,  "212",   "Malignant\n(Cancer)",      RED),
    (9.5,  "30",    "Measurements\nper Patient",YELLOW),
]:
    rect(sl, x, 1.3, 2.8, 1.8, RGBColor(0x2a, 0x2a, 0x4e))
    box(sl, val,   x+0.1, 1.4, 2.6, 0.9, font_size=36, bold=True,
        color=col, align=PP_ALIGN.CENTER)
    box(sl, label, x+0.1, 2.2, 2.6, 0.8, font_size=13,
        color=WHITE, align=PP_ALIGN.CENTER)

box(sl, "What Are the 30 Measurements?",
    0.5, 3.4, 12, 0.5, font_size=18, bold=True, color=YELLOW)

measurements = [
    ("Radius",    "How big the cell is"),
    ("Texture",   "How rough or smooth"),
    ("Perimeter", "Distance around the cell"),
    ("Area",      "Total size of the cell"),
    ("Concavity", "Number of dents in the cell"),
]
for i, (name, desc) in enumerate(measurements):
    x = 0.5 + (i % 3) * 4.2
    y = 4.0 + (i // 3) * 1.1
    rect(sl, x, y, 3.8, 0.9, RGBColor(0x2a, 0x2a, 0x4e))
    box(sl, f"{name}: {desc}", x+0.1, y+0.1, 3.6, 0.7,
        font_size=13, color=WHITE)

box(sl, "Each measurement is recorded 3 times: Mean · Standard Error · Worst  =  30 total",
    0.5, 6.8, 12, 0.5, font_size=13, color=LIGHT_GREY, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 4 — DATA CLEANING
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Step 2 — Cleaning the Data", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

box(sl, "Before teaching the computer we had to clean the data.\n"
        "Think of it like washing vegetables before cooking.",
    0.5, 1.1, 12, 0.7, font_size=16, color=LIGHT_GREY, italic=True)

steps = [
    ("1", "Removed the ID Column",
     "Every patient had a random ID number like 842302.\n"
     "This tells the computer nothing about cancer.\nWe deleted it."),
    ("2", "Checked for Missing Values",
     "Empty cells confuse the computer.\n"
     "We checked all 569 patients — no missing values found.\nThe data was complete."),
    ("3", "Scaled the Numbers",
     "Area = 2500  vs  Smoothness = 0.05\n"
     "The computer thinks bigger numbers are more important.\n"
     "Scaling makes all numbers the same size so they are treated fairly."),
]

for i, (num, title, desc) in enumerate(steps):
    x = 0.5 + i * 4.2
    rect(sl, x, 2.0, 3.8, 4.5, RGBColor(0x2a, 0x2a, 0x4e))
    rect(sl, x, 2.0, 3.8, 0.7, RED)
    box(sl, f"Fix {num}: {title}", x+0.1, 2.05, 3.6, 0.6,
        font_size=14, bold=True, color=WHITE)
    box(sl, desc, x+0.1, 2.8, 3.6, 3.5,
        font_size=13, color=WHITE)

box(sl, "Result: Clean, scaled data ready for the models",
    0.5, 6.7, 12, 0.5, font_size=14,
    color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — KEY FINDINGS FROM THE DATA
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Step 3 — What the Data Told Us", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

findings = [
    ("📏", "Cancer Cells Are Bigger",
     "Cancerous cells have larger radius, area,\nand perimeter than healthy cells."),
    ("🔴", "Cancer Cells Are More Irregular",
     "Concavity and concave points are\nmuch higher in cancerous cells."),
    ("🎯", "Strongest Clue",
     "concave_points_worst scored 0.79\n— the most powerful cancer indicator."),
    ("✅", "Clear Separation",
     "The two groups look very different.\nThe computer has clear patterns to learn."),
]

for i, (icon, title, desc) in enumerate(findings):
    x = 0.5 + (i % 2) * 6.4
    y = 1.3 + (i // 2) * 2.8
    rect(sl, x, y, 6.0, 2.4, RGBColor(0x2a, 0x2a, 0x4e))
    box(sl, f"{icon}  {title}", x+0.2, y+0.1, 5.6, 0.6,
        font_size=16, bold=True, color=YELLOW)
    box(sl, desc, x+0.2, y+0.7, 5.6, 1.5,
        font_size=14, color=WHITE)

img(sl, "chart3_top_features.png", 7.0, 1.2, 6.0, 5.8)


# ════════════════════════════════════════════════════════════
# SLIDE 6 — THE 6 MODELS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Step 4 — Teaching the Computer (6 Models)", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

box(sl, "Think of these as 6 different doctors — each with their own way of thinking:",
    0.5, 1.1, 12, 0.5, font_size=15, color=LIGHT_GREY, italic=True)

model_info = [
    ("Logistic Regression", "Draws a straight line to separate cancer from no cancer", "96.5%"),
    ("Decision Tree",       "Asks yes/no questions like a game of 20 questions",       "93.0%"),
    ("Random Forest",       "100 decision trees all voting together",                   "97.4%"),
    ("Gradient Boosting",   "Each tree learns from the mistakes of the last one",       "96.5%"),
    ("SVM",                 "Finds the best possible boundary between the two groups",  "97.4%"),
    ("KNN",                 "Looks at the 5 most similar patients and copies them",     "95.6%"),
]

for i, (name, desc, acc) in enumerate(model_info):
    x = 0.4 + (i % 2) * 6.4
    y = 1.8 + (i // 2) * 1.7
    rect(sl, x, y, 6.0, 1.5, RGBColor(0x2a, 0x2a, 0x4e))
    col = GREEN if float(acc[:-1]) >= 97 else (YELLOW if float(acc[:-1]) >= 96 else LIGHT_GREY)
    box(sl, name, x+0.15, y+0.1, 4.0, 0.5, font_size=14, bold=True, color=col)
    box(sl, desc, x+0.15, y+0.55, 4.5, 0.8, font_size=12, color=WHITE)
    box(sl, acc,  x+4.8,  y+0.3,  1.0, 0.7, font_size=18, bold=True,
        color=col, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 7 — RESULTS
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Step 5 — Results", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

try:
    with open("model_results.json") as f:
        results = json.load(f)
except:
    results = {}

headers = ["Model", "Accuracy", "Recall", "F1 Score"]
col_x   = [0.4, 4.5, 7.5, 10.3]
col_w   = [3.8, 2.7, 2.5, 2.7]

rect(sl, 0.4, 1.1, 12.5, 0.55, RED)
for hdr, x, w in zip(headers, col_x, col_w):
    box(sl, hdr, x, 1.15, w, 0.45, font_size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)

model_order = ["SVM", "Random Forest", "Logistic Regression",
               "Gradient Boosting", "KNN", "Decision Tree"]
for i, name in enumerate(model_order):
    y    = 1.8 + i * 0.75
    bg_c = RGBColor(0x2a, 0x2a, 0x4e) if i % 2 == 0 else RGBColor(0x22, 0x22, 0x3e)
    rect(sl, 0.4, y, 12.5, 0.72, bg_c)
    box(sl, name, col_x[0], y+0.1, col_w[0], 0.55,
        font_size=13, color=WHITE)
    if name in results:
        r = results[name]
        for val_key, x, w in zip(
            ["accuracy","recall","f1"], col_x[1:], col_w[1:]
        ):
            val = r[val_key]
            col = GREEN if val >= 97 else (YELLOW if val >= 93 else LIGHT_GREY)
            box(sl, f"{val:.1f}%", x, y+0.1, w, 0.55,
                font_size=13, bold=(val>=97), color=col,
                align=PP_ALIGN.CENTER)

rect(sl, 0.4, 6.4, 12.5, 0.7, RGBColor(0x1a, 0x4a, 0x1a))
box(sl, "🏆  Winner: SVM — 97.4% Accuracy · 97.5% Cross Validation · AUC 0.995",
    0.6, 6.5, 12, 0.55, font_size=15, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — CONFUSION MATRIX EXPLAINED
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "What Did the Model Get Right and Wrong?", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

img(sl, "chart6_confusion_matrix.png", 0.3, 1.1, 8.5, 4.5)

box(sl, "How to Read This:",
    9.2, 1.2, 3.8, 0.5, font_size=16, bold=True, color=YELLOW)

explanations = [
    ("✅ Top Left",  "Correctly said NO CANCER\n— and was right"),
    ("✅ Bot Right", "Correctly said CANCER\n— and was right"),
    ("⚠️ Top Right", "Said CANCER but was wrong\n— False Alarm"),
    ("❌ Bot Left",  "Said NO CANCER but missed it\n— Most Dangerous"),
]
for i, (label, desc) in enumerate(explanations):
    y = 1.9 + i * 1.3
    rect(sl, 9.2, y, 3.8, 1.2, RGBColor(0x2a, 0x2a, 0x4e))
    box(sl, label, 9.4, y+0.05, 3.4, 0.4, font_size=13, bold=True, color=YELLOW)
    box(sl, desc,  9.4, y+0.45, 3.4, 0.7, font_size=12, color=WHITE)

box(sl, "SVM missed only 3 cancer cases out of 114 patients tested",
    0.4, 5.9, 12.5, 0.5, font_size=14,
    color=GREEN, align=PP_ALIGN.CENTER, bold=True)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — THE WEB APP
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Step 6 — The Web App", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

box(sl, "We built a web app so that anyone can use the model — no coding needed.",
    0.5, 1.1, 12, 0.5, font_size=15, color=LIGHT_GREY, italic=True)

features = [
    ("🔬", "Detection Tab",
     "Enter 30 cell measurements.\nClick Run Detection.\nGet Cancer / No Cancer instantly."),
    ("📊", "Confidence Score",
     "Shows how sure the model is.\nExample: 94.2% confident — Malignant."),
    ("📈", "Suspicious Measurements",
     "Shows the top 5 measurements\nthat triggered the cancer flag."),
    ("🔄", "Model Switcher",
     "Choose between SVM,\nRandom Forest, or\nLogistic Regression."),
    ("📋", "Results Tab",
     "Full table of all 6 model scores.\nCharts and ROC curves included."),
    ("ℹ️", "About Tab",
     "Explains the dataset,\nwhat measurements mean,\nand all analysis charts."),
]

for i, (icon, title, desc) in enumerate(features):
    x = 0.4 + (i % 3) * 4.3
    y = 1.8 + (i // 3) * 2.5
    rect(sl, x, y, 4.0, 2.2, RGBColor(0x2a, 0x2a, 0x4e))
    box(sl, f"{icon}  {title}", x+0.15, y+0.1, 3.7, 0.55,
        font_size=14, bold=True, color=YELLOW)
    box(sl, desc, x+0.15, y+0.65, 3.7, 1.4, font_size=13, color=WHITE)

box(sl, "Run:  streamlit run breast_cancer_app.py  →  http://localhost:8503",
    0.4, 7.0, 12.5, 0.4, font_size=13,
    color=YELLOW, align=PP_ALIGN.CENTER, italic=True)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — SUMMARY
# ════════════════════════════════════════════════════════════
sl = add_slide(); bg(sl)
rect(sl, 0, 0, 13.33, 1.0, RED)
box(sl, "Summary — What We Built", 0.3, 0.1, 12, 0.8,
    font_size=28, bold=True, color=WHITE)

steps_summary = [
    ("1", "Explored the Data",     "569 patients · 30 measurements · 0 missing values"),
    ("2", "Cleaned the Data",       "Removed ID · Scaled all numbers fairly"),
    ("3", "Visualised the Data",    "4 charts showing cancer patterns clearly"),
    ("4", "Trained 6 Models",       "Logistic Regression · Decision Tree · Random Forest · SVM · KNN · Gradient Boosting"),
    ("5", "Evaluated the Models",   "SVM won — 97.5% cross validation · AUC 0.995 · Only 3 missed cancers"),
    ("6", "Built a Web App",        "Anyone can enter measurements and get a result instantly"),
]

for i, (num, title, desc) in enumerate(steps_summary):
    y = 1.2 + i * 0.95
    rect(sl, 0.4, y, 12.5, 0.88, RGBColor(0x2a, 0x2a, 0x4e))
    rect(sl, 0.4, y, 0.7,  0.88, RED)
    box(sl, num, 0.4, y+0.15, 0.7, 0.55,
        font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(sl, title, 1.3, y+0.05, 3.5, 0.45,
        font_size=14, bold=True, color=YELLOW)
    box(sl, desc,  4.9, y+0.1,  7.8, 0.7,
        font_size=13, color=WHITE)

rect(sl, 0.4, 7.0, 12.5, 0.4, RGBColor(0x1a, 0x4a, 0x1a))
box(sl, "Final Result: A machine learning system that detects breast cancer with 97.5% accuracy",
    0.6, 7.05, 12, 0.35, font_size=13, bold=True,
    color=GREEN, align=PP_ALIGN.CENTER)


# ── SAVE ─────────────────────────────────────────────────────
prs.save("Breast_Cancer_Detection_Presentation.pptx")
print("Presentation saved: Breast_Cancer_Detection_Presentation.pptx")
print(f"Total slides: {len(prs.slides)}")
